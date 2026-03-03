from typing import Any

from pptx import Presentation
from pptx.util import Inches

from app.services.chunking.base import ChunkResult, ChunkingStrategy
from app.utils.text import clean_text, count_tokens

MIN_CHUNK_TOKENS = 30


class PptxChunker(ChunkingStrategy):
    """Chunks PowerPoint files: extract slides with titles, bullets, notes."""

    async def chunk(self, content: Any, metadata: dict[str, Any]) -> list[ChunkResult]:
        """Content should be a file path to the .pptx file."""
        filepath: str = content
        doc_title = metadata.get("title", "")
        prs = Presentation(filepath)

        slides_data = _extract_slides(prs)
        chunks: list[ChunkResult] = []

        for i, slide in enumerate(slides_data):
            text = slide["text"]
            if count_tokens(text) < MIN_CHUNK_TOKENS:
                continue

            header = self._build_context_header(
                doc_title,
                section=slide["title"],
                prev_topic=chunks[-1].topic_label if chunks else None,
            )

            full_content = f"[{header}]\n\n"
            if slide["title"]:
                full_content += f"# {slide['title']}\n\n"
            full_content += text
            if slide["notes"]:
                full_content += f"\n\nNotizen: {slide['notes']}"

            chunks.append(ChunkResult(
                content=clean_text(full_content),
                chunk_index=len(chunks),
                chunk_type=_detect_slide_type(text),
                section_heading=slide["title"],
                page_number=i + 1,
                metadata={
                    "slide_number": i + 1,
                    "has_images": slide["has_images"],
                },
            ))

        return chunks


def _extract_slides(prs: Presentation) -> list[dict]:
    """Extract structured data from each slide."""
    slides = []
    for slide in prs.slides:
        title = ""
        body_parts: list[str] = []
        has_images = False

        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    title = shape.text.strip()
                else:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            level = para.level
                            prefix = "  " * level + "- " if level > 0 else ""
                            body_parts.append(f"{prefix}{text}")

            if shape.shape_type == 13:  # Picture
                has_images = True

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "title": title,
            "text": "\n".join(body_parts),
            "notes": notes,
            "has_images": has_images,
        })

    return slides


def _detect_slide_type(text: str) -> str:
    lower = text.lower()
    if any(w in lower for w in ("definition", "def:", "definiert")):
        return "definition"
    if any(w in lower for w in ("beispiel", "example")):
        return "example"
    if any(w in lower for w in ("agenda", "gliederung", "inhalt", "übersicht")):
        return "overview"
    return "theory"
